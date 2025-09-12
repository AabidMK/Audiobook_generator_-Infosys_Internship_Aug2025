from pptx import Presentation
def extract_ppt_text(ppt_path):
    prs = Presentation(ppt_path)
    text = []

    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text.append(shape.text)

    return "\n".join(text)