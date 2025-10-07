from pptx import Presentation
def extract_ppt_text(ppt_path):
    prs = Presentation(ppt_path)
    text = ""
    
    for slide_num, slide in enumerate(prs.slides, start=1):
        text += f"\n## Slide {slide_num}\n"
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                text += shape.text.strip() + "\n"
    return text