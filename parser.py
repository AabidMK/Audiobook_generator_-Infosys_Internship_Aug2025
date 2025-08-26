import os
import pdfplumber
import docx
from PIL import Image
import pytesseract
from pptx import Presentation
pytesseract.pytesseract.tesseract_cmd = r"C:\Program Files\Tesseract-OCR\tesseract.exe"

def extract_pdf_text(pdf_path):
    text = ""
    with pdfplumber.open(pdf_path) as pdf:
        for page in pdf.pages:
            content = page.extract_text()
            if content:
                text += content + "\n\n"
    return text

def extract_ppt_text(ppt_path):
    prs = Presentation(ppt_path)
    text = []

    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text.append(shape.text)

    return "\n".join(text)

def extract_docx_text(docx_path):
    doc = docx.Document(docx_path)
    text = "\n".join([para.text for para in doc.paragraphs if para.text.strip()])

    for table in doc.tables:
        for row in table.rows:
            row_data = [cell.text.strip() for cell in row.cells]
            text += "\n" + " | ".join(row_data)
    return text


def extract_image_text(image_path, lang="eng"):
    try:
        img = Image.open(image_path)
        text = pytesseract.image_to_string(img, lang=lang)
        return text.strip()
    except Exception as e:
        return f"Error: {e}"


def write_to_markdown(content, md_path):
    with open(md_path, "w", encoding="utf-8") as f:
        f.write(content)


def extract_to_markdown(file_path, md_path="output.md"):
    ext = os.path.splitext(file_path)[1].lower()

    if ext == ".pdf":
        content = extract_pdf_text(file_path)
    elif ext == ".docx":
        content = extract_docx_text(file_path)
    elif ext in [".png", ".jpg", ".jpeg", ".bmp", ".tiff"]:
        content = extract_image_text(file_path)
    elif ext==".pptx":
        content= extract_ppt_text(file_path)    
    else:
        raise ValueError("Unsupported file format. Please provide a PDF, DOCX, or image.")

    write_to_markdown(content, md_path)
